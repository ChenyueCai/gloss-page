#!/usr/bin/env python3
"""Re-tone slides in the DECK, not in the exported bitmap.

Doing it here keeps the labels as real vector text, so they export crisp at any
resolution. Post-processing a render can only guess which pixels are type, and
that guessing is what made the earlier exports look mottled.

Rule: the slide ground becomes the page ground and type on it becomes light.
Type sitting inside one of the pale panels stays dark, because those panels keep
their own colour and dark-on-pale is what reads there.
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

BG   = RGBColor(0x0d, 0x12, 0x0e)     # --bg
INK  = RGBColor(0xe9, 0xee, 0xea)     # --ink
PALE = {"A7D1E4", "EECAE8"}           # panels that keep dark text on top
MARKERS = set("ABCDEF1234567")        # the circled step markers


IN = 914400
ORANGE = RGBColor(0xf6, 0x97, 0x46)   # --accent-2
RULE   = RGBColor(0x7b, 0x85, 0x7d)   # --ink-3

# ---- diagram palette ------------------------------------------------------
# Re-toning the type alone left the drawing in the paper's colours: navy rules,
# pale pink and pale blue panels, a magenta block, light-blue arrows. Against
# --bg that reads as a slide someone pasted onto the page. These map the deck's
# drawn shapes onto the page's own roles - green carries structure, orange is
# reserved for emphasis - and touch nothing that is a picture.
ACCENT   = RGBColor(0x5d, 0xcb, 0x81)   # --accent: flow, structure
DEEP     = RGBColor(0x2a, 0x97, 0x54)   # --accent-deep: filled process boxes
ACCENT_I = RGBColor(0x04, 0x11, 0x07)   # --accent-ink: type on an orange fill
PANEL    = RGBColor(0x1c, 0x22, 0x1e)   # --panel
LINE     = RGBColor(0x2c, 0x34, 0x2e)   # --line
PANEL_G  = RGBColor(0x15, 0x2a, 0x1d)   # automatic-texturing panel, green-tinted
PANEL_O  = RGBColor(0x2d, 0x21, 0x18)   # interactive-texturing panel, warm

# source hex -> (fill, line) once the shape is not a panel
DIAGRAM = {
    "A7D1E4": (ACCENT,  ACCENT),
    "ABABAB": (ACCENT,  ACCENT),    # slide 2 grey arrows: same role, same colour    # flow arrows and their connectors
    "EECAE8": (ORANGE,  ORANGE),    # the loss arrows on the data slide
    "E632BB": (ORANGE,  None),      # the GLOSS block
    "D778F5": (PANEL_O, None),      # the GLOSS network box
    "F2F215": (ORANGE,  None),      # the batch multi-attention pill
    "17649B": (DEEP,    LINE),      # filled process boxes / their outlines
    "4472C4": (DEEP,    DEEP),
    "EBECDE": (PANEL,   LINE),      # cream label boxes
    "ABABAB": (RULE,    RULE),      # the mask polygons over the albedo tiles
}
PANELS = {"A7D1E4": PANEL_G, "EECAE8": PANEL_O}
DIAGRAM_SLIDES = {2, 20}

# The rest of slide 2 takes its colour from the PowerPoint shape STYLE - a
# theme reference rather than an explicit RGB - so nothing above matches it and
# the hex is not readable from spPr. Those shapes are reached by what they are
# and what they say, which is stable across re-exports of the deck.
BY_SHAPE = {
    "TRAPEZOID":   (DEEP,  INK),      # the E / D encoder and decoder blocks
    "LEFT_BRACE":  (None,  None),     # brackets: outline only, handled below
    "RIGHT_BRACE": (None,  None),
}
BY_LABEL = {
    "ControlNet":    (DEEP,    INK),
    "RGB\nto PBR\nNetwork": (DEEP, INK),
    "GLOSS\nNetwork": (PANEL_O, INK),
    "random masks":  (PANEL,   INK),
}
BRACES = {"LEFT_BRACE", "RIGHT_BRACE"}
ON_ORANGE = {"E632BB", "F2F215"}    # fills whose own type has to go dark


def _hex(colour):
    try:
        return str(colour.rgb).upper()
    except Exception:
        return None


def shape_fill_hex(sh):
    try:
        return _hex(sh.fill.fore_color) if sh.fill.type == 1 else None
    except Exception:
        return None


def shape_line_hex(sh):
    """Read the outline colour straight out of the XML.

    Going through sh.line.color would be the obvious way, but python-pptx sets
    the line fill to solid as a side effect of handing back the ColorFormat.
    Merely asking a label what colour its border is would therefore give it
    one, and the first pass drew a box round every label on the slide.
    """
    ln = _ln(sh)
    if ln is None:
        return None
    solid = ln.find(qn('a:solidFill'))
    if solid is None:
        return None
    clr = solid.find(qn('a:srgbClr'))
    return clr.get('val').upper() if clr is not None else None


def _ln(sh):
    try:
        return sh._element.spPr.find(qn('a:ln'))
    except Exception:
        return None


def set_fill(sh, colour):
    try:
        sh.fill.solid()
        sh.fill.fore_color.rgb = colour
        sh.shadow.inherit = False
    except Exception:
        pass


def line_visible(sh):
    """Only shapes that already draw a line may have one recoloured: an
    explicit solid outline, or a connector, whose whole substance is its line."""
    if "LINE" in str(sh.shape_type) or "Connector" in sh.name:
        return True
    ln = _ln(sh)
    return ln is not None and ln.find(qn('a:solidFill')) is not None


def set_line(sh, colour, force=False):
    if not force and not line_visible(sh):
        return False
    try:
        sh.line.color.rgb = colour
        return True
    except Exception:
        return False


def set_text(sh, colour):
    if not sh.has_text_frame:
        return 0
    n = 0
    for para in sh.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = colour
            n += 1
    return n


def is_panel(sh):
    """The two big rounded rectangles the target meshes sit on. They are the
    only pale shapes wide enough to be a ground rather than a mark."""
    try:
        return (str(sh.auto_shape_type).startswith("ROUNDED_RECTANGLE")
                and sh.width / IN > 3.0 and sh.height / IN > 2.0)
    except Exception:
        return False


# A few blocks take their outline from the theme rather than explicit RGB, so
# the hex map cannot see them. The GLOSS block is the emphasis of the whole
# figure and was drawing a pink theme border over its warm fill.
OUTLINE_BY_LABEL = {
    "GLOSS\nNetwork": ORANGE,
}


def accent_theme_connectors(shapes, seen):
    """Connector LINEs whose colour comes from the theme, not explicit RGB.

    Slide 2 draws 13 of these and they stay the original grey while every
    hex-mapped arrow around them turns accent green, so the figure ends up with
    two different arrow colours doing one job. The hex map cannot see them for
    the same reason it could not see the GLOSS border.
    """
    for sh in shapes:
        if getattr(sh, "shapes", None):
            accent_theme_connectors(sh.shapes, seen)
            continue
        if not str(sh.shape_type).startswith("LINE"):
            continue
        if shape_line_hex(sh) is not None:
            continue                      # explicit colour: the hex map owns it
        set_line(sh, ACCENT, force=True)
        seen["line"] += 1


def apply_outlines(shapes, seen):
    for sh in shapes:
        if getattr(sh, "shapes", None):
            apply_outlines(sh.shapes, seen)
            continue
        if not sh.has_text_frame:
            continue
        want = OUTLINE_BY_LABEL.get(sh.text_frame.text.strip())
        if want is None:
            continue
        set_line(sh, want, force=True)   # force: there is no visible line to detect
        try:
            sh.line.width = Pt(1.75)
        except Exception:
            pass
        seen["line"] += 1


def recolour_diagram(shapes, seen):
    """Walk the drawn shapes and move them onto the page palette. Pictures are
    never touched: the renders, patch grids and camera icons are content."""
    for sh in shapes:
        t = str(sh.shape_type)
        if "GROUP" in t:
            recolour_diagram(sh.shapes, seen)
            continue
        if "PICTURE" in t:
            continue
        st = ""
        try:
            st = str(sh.auto_shape_type).split()[0]
        except Exception:
            pass
        label = ""
        try:
            if sh.has_text_frame:
                label = sh.text_frame.text.strip()
        except Exception:
            pass
        if st in BRACES:
            # a brace draws nothing but its outline, and that outline comes
            # from the style, so it has to be forced rather than matched
            if set_line(sh, seen.get("flow", RULE), force=True):
                seen["line"] += 1
            continue
        styled = BY_SHAPE.get(st) or BY_LABEL.get(label)
        if styled and styled[0] is not None:
            set_fill(sh, styled[0])
            set_text(sh, styled[1])
            seen["fill"] += 1
            continue
        f, l = shape_fill_hex(sh), shape_line_hex(sh)
        if f in PANELS and is_panel(sh):
            set_fill(sh, PANELS[f])
            seen["panel"] += 1
            continue
        if f in DIAGRAM:
            fill, line = DIAGRAM[f]
            set_fill(sh, fill)
            if line is not None:
                set_line(sh, line)
            if f in ON_ORANGE:
                set_text(sh, ACCENT_I)
                seen["on_orange"].append(sh)
            seen["fill"] += 1
        elif l in DIAGRAM:
            if set_line(sh, DIAGRAM[l][1] or DIAGRAM[l][0]):
                seen["line"] += 1
        elif "LINE" in t or "Connector" in sh.name:
            # the plain rules that separate the three columns are theme-coloured
            # rather than explicit navy, so they never match the map above
            if set_line(sh, seen.get("flow", RULE)):
                seen["line"] += 1


def drop_bottom_strip(slide, min_w_in=8.0, max_h_in=2.2, min_top_in=5.5):
    """Slide 2 carries a full-width dark band across its foot that belongs to
    the talk layout, not the figure. Remove it rather than cropping, so the
    figure keeps its own margins."""
    gone = 0
    for sh in list(slide.shapes):
        try:
            w, h, t = sh.width / IN, sh.height / IN, sh.top / IN
        except TypeError:
            continue
        if w >= min_w_in and h <= max_h_in and t >= min_top_in:
            sh._element.getparent().remove(sh._element)
            gone += 1
    return gone


def outline(slide, x, y, w, h, colour, pt):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(x * IN)), Emu(int(y * IN)),
                                 Emu(int(w * IN)), Emu(int(h * IN)))
    box.fill.background()
    box.line.color.rgb = colour
    box.line.width = Pt(pt)
    box.shadow.inherit = False
    return box


def unclip_row_labels(slide):
    """The rotated row labels on slide 21 sit in boxes a shade too short for
    their type, so LibreOffice wraps them ("Albed / o"). Turn wrapping off and
    give each box room; the rotation keeps them in place."""
    fixed = 0
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        if sh.text_frame.text.strip() not in ("Albedo", "Normals", "Output"):
            continue
        # Widening the box is the wrong lever: these sit at a negative x, so a
        # wider box runs off the slide edge and over the first grid column.
        # Shrink the type instead and the original geometry still holds.
        try:
            sh.text_frame.word_wrap = False
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    pts = run.font.size.pt if run.font.size else 12.0
                    run.font.size = Pt(round(pts * 0.82, 1))
            fixed += 1
        except Exception:
            pass
    return fixed


def annotate_attention(slide):
    """Two examples sit side by side with nothing between them, and the target
    patch - the last column of each - is not called out. Add both."""
    left  = (0.28, 8.78)      # measured from the picture extents
    right = (8.90, 17.40)
    top, bot = 0.55, 3.74
    cols = 8
    # separator down the gap between the two examples
    gap = (left[1] + right[0]) / 2
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int((gap - 0.008) * IN)),
                                  Emu(int((top - 0.16) * IN)), Emu(int(0.016 * IN)),
                                  Emu(int((bot - top + 0.32) * IN)))
    rule.fill.solid(); rule.fill.fore_color.rgb = RULE
    rule.line.fill.background(); rule.shadow.inherit = False
    # the target patch is the final column of each example
    for x0, x1 in (left, right):
        colw = (x1 - x0) / cols
        outline(slide, x1 - colw, top, colw, bot - top, ORANGE, 2.0)


def panel_boxes(slide):
    out = []
    for sh in slide.shapes:
        try:
            f = sh.fill
            if f.type == 1 and str(f.fore_color.rgb).upper() in PALE:
                out.append((sh.left, sh.top, sh.left + sh.width, sh.top + sh.height))
        except Exception:
            pass
    return out

def light_filled(sh):
    """True when the shape carries its own pale fill - a white marker disc, say.
    Its text has to stay dark or it disappears into the fill."""
    try:
        f = sh.fill
        if f.type != 1:
            return False
        c = f.fore_color.rgb
        return (c[0] + c[1] + c[2]) / 3 > 140
    except Exception:
        return False


def inside(sh, boxes):
    try:
        cx, cy = sh.left + sh.width // 2, sh.top + sh.height // 2
    except TypeError:
        return False
    return any(x0 <= cx <= x1 and y0 <= cy <= y1 for x0, y0, x1, y1 in boxes)

def retone(path_in, path_out, slide_numbers):
    prs = Presentation(path_in)
    touched = 0
    for n in slide_numbers:
        slide = prs.slides[n - 1]
        fill = slide.background.fill
        fill.solid(); fill.fore_color.rgb = BG
        if n == 2:
            print("  slide 2: removed", drop_bottom_strip(slide), "layout strip(s)")
        if n == 21:
            annotate_attention(slide)
            print("  slide 21: unclipped", unclip_row_labels(slide), "row label(s)")
            print("  slide 21: separator + target-patch boxes added")
        # The drawing moves onto the page palette first. That also empties
        # panel_boxes() below - once the panels are dark there is no pale
        # ground left for dark type to sit on, so every label and every step
        # disc on the slide flips light by the existing rules.
        if n in DIAGRAM_SLIDES:
            seen = {"panel": 0, "fill": 0, "line": 0, "on_orange": []}
            seen["flow"] = ACCENT if n == 2 else RULE
            recolour_diagram(slide.shapes, seen)
            apply_outlines(slide.shapes, seen)
            accent_theme_connectors(slide.shapes, seen)
            print("  slide %d: %d panel(s), %d filled shape(s), %d rule(s) recoloured"
                  % (n, seen["panel"], seen["fill"], seen["line"]))
        boxes = panel_boxes(slide)
        for sh in slide.shapes:
            if not sh.has_text_frame or not sh.text_frame.text.strip():
                continue
            # A step disc is an oval. Without that test the rule also catches
            # the E and D encoder blocks on slide 2, whose whole label is a
            # single letter, and repaints them as if they were step numbers.
            is_oval = False
            try:
                is_oval = str(sh.auto_shape_type).startswith("OVAL")
            except Exception:
                pass
            marker = is_oval and sh.text_frame.text.strip() in MARKERS
            if marker and not inside(sh, boxes):
                # Step discs are white with dark type. Flipping only the type
                # makes them white-on-white, so flip the whole disc instead:
                # dark ground, ink outline, ink numeral.
                try:
                    sh.fill.solid(); sh.fill.fore_color.rgb = BG
                    sh.line.color.rgb = INK; sh.line.width = Pt(1.0)
                    sh.shadow.inherit = False
                except Exception:
                    pass
                for para in sh.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = INK
                        touched += 1
                continue
            if inside(sh, boxes) or light_filled(sh):
                continue          # dark type on a pale panel or disc: leave it
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = INK
                    touched += 1

    prs.save(path_out)
    print(f"recoloured {touched} text runs across slides {slide_numbers}")

if __name__ == "__main__":
    retone("figs.pptx", "figs_dark.pptx", [2, 20, 21])
